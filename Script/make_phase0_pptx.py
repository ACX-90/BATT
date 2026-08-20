"""第 0 期实验汇报：白底、16 号字、每页约一半留白，便于贴进模板。"""

from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

OUT = "Doc/04-c-第0期实验汇报.pptx"
FONT = "Microsoft YaHei"
INK = RGBColor(0x22, 0x22, 0x22)
MUTED = RGBColor(0x55, 0x55, 0x55)
LINE = RGBColor(0xB0, 0xB0, 0xB0)
HEAD = RGBColor(0xF0, 0xF0, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NSMAP_A = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}


def _ea(run) -> None:
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        rPr.append(
            parse_xml(
                '<a:ea xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
                f' typeface="{FONT}"/>'
            )
        )
    else:
        ea.set("typeface", FONT)


def style_run(run, size=16, bold=False, color=INK) -> None:
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT
    _ea(run)


def add_title(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.5))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    style_run(run, 18, True)


def add_body(slide, lines: list[str], *, y=1.05, w=6.3, h=5.8, size=16) -> None:
    """Left column ~ half width. Remaining area is for figures."""
    box = slide.shapes.add_textbox(Inches(0.6), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        p.line_spacing = 1.15
        run = p.add_run()
        run.text = line
        style_run(run, size, False)


def add_table(slide, rows: list[list[str]], *, y=1.05, w=12.1, col_w=None) -> None:
    n_row, n_col = len(rows), len(rows[0])
    row_h = Inches(0.40)
    tbl_h = row_h * n_row
    shape = slide.shapes.add_table(n_row, n_col, Inches(0.6), Inches(y), Inches(w), tbl_h)
    table = shape.table
    table.style = "Table Grid"
    if col_w:
        for i, cw in enumerate(col_w):
            table.columns[i].width = Inches(cw)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = ""
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = HEAD if r == 0 else WHITE
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.08)
            tf.margin_right = Inches(0.08)
            tf.margin_top = Inches(0.05)
            tf.margin_bottom = Inches(0.05)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = val
            style_run(run, 16, bold=(r == 0))


def blank_slide(prs: Presentation):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def build() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1 title
    s = blank_slide(prs)
    box = s.shapes.add_textbox(Inches(0.6), Inches(2.15), Inches(10), Inches(0.5))
    tf = box.text_frame
    tf.margin_left = tf.margin_top = 0
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "第 0 期实验汇报"
    style_run(run, 22, True)
    add_body(
        s,
        [
            "100 Ah NMC 模板 · 一阶等效电路 · 合成数据对照",
            "英飞凌 demo 短窗口不够当过关，我们补 SOC 与电阻怎么拆",
        ],
        y=2.8,
        w=10,
        h=1.6,
        size=16,
    )

    # 2 结论
    s = blank_slide(prs)
    add_title(s, "结论")
    add_body(
        s,
        [
            "第 0 期已经做完。A 到 H 都有正式数字，可以进第 1 期。",
            "这期不是再做一个更炫的小网。",
            "英飞凌 demo 放 300 秒可以报 SOC 1% 以内；时间一拉长，容量错和电流零偏会把安时拉开。",
            "要先证明：电阻表该不该改、该怎么改。不能只看电压变好，更不能只看五分钟的 SOC。",
        ],
    )

    # 3 demo
    s = blank_slide(prs)
    add_title(s, "英飞凌 demo 不够好，窗口还太短")
    add_body(
        s,
        [
            "现场结构：3×8×2 小网，电阻 = 规格书点 + 网络残差。电压偏过 10 mV 就按拍反传。",
            "对外放了大约 300 秒，SOC 误差报 1% 以内。",
            "300 秒在 1C 上只走约 8 个百分点 SOC。容量错 5%，或零偏 5 A，积分都只有约 0.4 个百分点，还在 1% 里。",
            "这不是估准了，是安时误差还没长出来。车上是快充几十分钟、停车几小时，不是 300 秒展台。",
        ],
    )

    # 4 300s table
    s = blank_slide(prs)
    add_title(s, "300 秒结构上验不出积分误差")
    add_table(
        s,
        [
            ["错配", "300 秒、1C 上安时会偏", "进不进 1%"],
            ["容量打九五折", "约 0.44 个百分点", "进"],
            ["5 A 电流零偏", "约 0.42 个百分点", "进"],
            ["容量错 10%", "约 0.93 个百分点", "刚进"],
        ],
        y=1.1,
        w=8.5,
        col_w=[2.6, 3.5, 2.4],
    )
    add_body(
        s,
        ["要让 5% 容量错配顶破 1%，净安时得超过约 19 Ah：1C 大约 11 分钟，0.3C 大约 38 分钟。"],
        y=3.0,
        w=8.5,
        h=1.4,
    )

    # 5 小时级
    s = blank_slide(prs)
    add_title(s, "同一套错配，拉到两小时")
    add_table(
        s,
        [
            ["错配", "安时终点", "卡尔曼终点", "开环电压"],
            ["容量 95 Ah", "−3.16 pp", "+0.04 pp", "21.9 mV"],
            ["容量 105 Ah", "+2.86 pp", "+0.08 pp", "9.5 mV"],
            ["5 A 零偏", "−10.0 pp", "+1.40 pp", "143 mV"],
        ],
        y=1.1,
        w=10.5,
        col_w=[2.5, 2.6, 2.7, 2.7],
    )
    add_body(
        s,
        [
            "工况：0.3C 放电 2 小时。仿真仍是 100 Ah。",
            "安时跑飞；卡尔曼用电压把 SOC 钉回去。若再把电压误差写进电阻，SOC 和内阻一起飞。",
        ],
        y=3.1,
        w=10.5,
        h=1.6,
    )

    # 6 +3%
    s = blank_slide(prs)
    add_title(s, "跑飞之后，会停在真值 +3% 附近")
    add_body(
        s,
        [
            "开环安时不会自己停。能稳住，说明 demo 用电压在往回拉。",
            "停在 +3% 而不是 0，是 10 mV 死区折到中段开路电压斜率上。",
            "中段大约每 3 mV 对应 1% SOC：10 mV ÷ 3 mV ≈ 3%。",
            "这不是容量刚好错了 3%。那种误差会随安时继续长。",
            "我们的卡尔曼没有这道 SOC 死区。同样打到 105 Ah，安时 +2.86 pp，滤波终点只有 +0.08 pp。",
        ],
    )

    # 7 解决什么
    s = blank_slide(prs)
    add_title(s, "我们要补的不是小网，是这三件事")
    add_body(
        s,
        [
            "1. 时间拉长以后，SOC 还能不能钉住。用卡尔曼吃端电压，不要让电阻头去跟斜坡。",
            "2. 电压对不上时，到底该不该改电阻表。四种「新」不能用同一套验收。",
            "3. 不要照搬「过 10 mV 就每 0.1 秒反传」。那是没回放时的展台折中。",
            "一律再训练，新年份电压往往会降，旧工况的表却被改坏。",
        ],
    )

    # 8 设置
    s = blank_slide(prs)
    add_title(s, "实验怎么设")
    add_body(
        s,
        [
            "对象：100 Ah NMC 模板，不是某一款商用电芯。数字用来比方法。",
            "电路：一阶 Thevenin。网络只出 R0、R1，电容钉死。",
            "增量损失：开环电压误差。滤波把电压贴回去之后的残差不拿来训练。",
            "训练网格大约 10 分钟一条，带边沿和回弹。小时级另做，不进训练集。",
        ],
    )

    # 9 五档
    s = blank_slide(prs)
    add_title(s, "每种错配都用同一套五档")
    add_table(
        s,
        [
            ["档", "做什么", "适合什么时候"],
            ["冻结", "权重完全不动", "其实没有新东西"],
            ["合集重训", "旧网格和新数据一起再训", "旧数据还在、想重拟合"],
            ["回放", "新年份为主，抽一部分旧轨迹", "补一块新区域、保住旧区域"],
            ["只微调", "只扫新年份", "容易把旧表改坏"],
            ["缩放", "只学两个正数乘子 k0、k1", "整张表幅度变了、形状还在"],
        ],
        y=1.1,
        w=12.0,
        col_w=[2.2, 4.4, 5.4],
    )

    # 10 A
    s = blank_slide(prs)
    add_title(s, "A. 整张电阻 ×1.15")
    add_table(
        s,
        [
            ["档", "新年份", "旧网格", "乘子"],
            ["冻结", "23 mV", "4 mV", "—"],
            ["缩放", "约 7 mV", "22 mV", "约 1.20，两通道齐"],
            ["只微调", "6 mV，最好", "18 mV", "两通道不齐"],
        ],
        y=1.1,
        w=10.5,
        col_w=[2.2, 2.6, 2.4, 3.3],
    )
    add_body(
        s,
        ["结论：整体涨阻，只动两个乘子。旧网格变差是缺寿命维的正常结果，不是遗忘失败。"],
        y=3.15,
        w=10.5,
        h=1.3,
    )

    # 11 B C
    s = blank_slide(prs)
    add_title(s, "B. 填洞  −10 °C　　C. 同分布再贴")
    add_body(
        s,
        [
            "B 舰队没见过 −10 °C。冻结：新 10 mV，旧 7.5 mV。回放 / 重训把新温区降到旧集量级。缩放两个乘子朝反方向走。",
            "结论：缺温区，用回放或重训，不要缩放。",
            "C 舰队其实已经见过这一档。冻结新 6 mV、旧 4 mV。只微调新年份削了零点几毫伏，旧集抬约 10%。",
            "结论：同分布再贴，正确动作是不更新。只报新年份电压会假装成功。",
        ],
        w=7.2,
    )

    # 12 E F
    s = blank_slide(prs)
    add_title(s, "E. 这一趟表偏　　F. 小时级门控")
    add_body(
        s,
        [
            "E 仿真仍是新电芯，只把表上的 R0 乘 1.2。关掉滤波慢偏置，休息 SOC 被拉约 0.5 pp；打开后偏置约 −160 µΩ，SOC 回来。不要当成老化去改整张表。",
            "F 健康长恒流：开环 1.2 mV，门控拒绝。零偏 5 A、容量打九五折：门控会通过，但滤波电压仍亚毫伏，安时已经漂了。",
            "结论：门控通过 ≠ 该改电阻。长恒流加停车不是增量主集。",
        ],
        w=7.4,
    )

    # 13 G H
    s = blank_slide(prs)
    add_title(s, "G. 寿命因子 0.90　　H. 真值是二阶")
    add_body(
        s,
        [
            "G 电阻、电容按寿命因子乘，开路电压和容量先不动。R0 大约 ×1.10，R1 大约 ×1.15。缩放两个乘子分别对着这两个数，不相等不是拆错通道。",
            "结论：缺寿命维的涨阻，仍走缩放。不要和容量掉、二阶电路第一次叠在一起。",
            "H 生成器多一条约 90 秒的慢极化，估计器不读。开环从 1 mV 到 12 mV，回弹后段留下 7–15 mV 同号慢尾巴。",
            "结论：这是缺一条时间常数，不是表幅度错了。不要用二阶电压去增量一阶网络。",
        ],
        w=7.4,
    )

    # 14 规则
    s = blank_slide(prs)
    add_title(s, "先看开环电压，再决定动哪一层")
    add_table(
        s,
        [
            ["你看到的", "该做什么", "不该做什么"],
            ["几乎没高", "冻结", "再训练"],
            ["各温区边沿都大，像电流×多出来的电阻", "只学两个乘子", "解冻整张网"],
            ["只在没覆盖的温度大", "回放或重训", "缩放"],
            ["回弹后段同号慢尾巴", "允许残差，不增量", "写成内阻"],
            ["开环大、滤波健康、安时在漂", "先查分流器和容量", "拆 R0 / R1"],
            ["只有这一趟边沿大", "滤波慢偏置", "立刻改网络"],
        ],
        y=1.05,
        w=12.1,
        col_w=[4.5, 3.8, 3.8],
    )

    # 15 四个问题
    s = blank_slide(prs)
    add_title(s, "事先四个问题")
    add_body(
        s,
        [
            "同分布再训练，会不会假装成功？会。只报新年份电压，微调看起来赢了。",
            "激励不足时电压变好，会不会把电阻拆错？会。现有门控挡不住零偏和容量错。过了门也不许拆通道。",
            "测量列相对真值列偏多少？墙仍是电阻幅度。底板差约 3.6 mV，主因是训次没对齐。这题只答了一半。",
            "开环残差长什么样？一阶看二阶真值，回弹后段是第二指数，大约 7–15 mV。不是增量对象。",
        ],
        w=7.4,
    )

    # 16 完成度
    s = blank_slide(prs)
    add_title(s, "完成度，和下一期")
    add_body(
        s,
        [
            "按计划，A 到 H 已齐。可以进入第 1 期。",
            "这期故意没做：不换循环骨干，不把网络塞进卡尔曼，不升二阶，不写 MCU 固件。",
            "记下了但还没补：测量列 / 填洞舰队只训了 100 轮；小时级没有强行微调。不影响结题口径。",
            "下一期：滑窗更新（没边沿、没回弹就不走一步）；每芯一个很小的残差头，纯涨阻时不该赢过两个乘子。",
        ],
    )

    prs.save(OUT)
    print(OUT, "slides", len(prs.slides))


if __name__ == "__main__":
    build()
